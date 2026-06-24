from pathlib import Path
from pptx import Presentation
from app.rag.parsers.base import ParsedChunk


class PptxParser:
    def parse(self, path: Path) -> list[ParsedChunk]:
        prs = Presentation(str(path))
        chunks = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.extend(p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip())
            if texts:
                chunks.append(ParsedChunk("\n".join(texts), {"slide": i + 1, "total_slides": len(prs.slides)}))
        return chunks
