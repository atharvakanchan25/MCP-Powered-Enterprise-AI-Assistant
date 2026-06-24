import pytest
from pathlib import Path
from app.rag.parsers import parse_document, SUPPORTED_EXTENSIONS


def test_txt_parser(tmp_path):
    f = tmp_path / "sample.txt"
    f.write_text("Hello world\nSecond line", encoding="utf-8")
    chunks = parse_document(f)
    assert len(chunks) == 1
    assert "Hello world" in chunks[0].text


def test_pdf_parser(tmp_path):
    # Build a minimal PDF manually (1-page text PDF)
    pdf_content = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R"
        b"/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj\n"
        b"4 0 obj<</Length 44>>\nstream\nBT /F1 12 Tf 100 700 Td (Test PDF) Tj ET\nendstream\nendobj\n"
        b"5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj\n"
        b"xref\n0 6\n0000000000 65535 f\n"
        b"trailer<</Size 6/Root 1 0 R>>\nstartxref\n0\n%%EOF"
    )
    f = tmp_path / "sample.pdf"
    f.write_bytes(pdf_content)
    # pypdf may or may not extract text from this minimal PDF; just ensure no crash
    chunks = parse_document(f)
    assert isinstance(chunks, list)


def test_docx_parser(tmp_path):
    from docx import Document
    doc = Document()
    doc.add_paragraph("First paragraph of the document.")
    doc.add_paragraph("Second paragraph with more content.")
    path = tmp_path / "sample.docx"
    doc.save(str(path))
    chunks = parse_document(path)
    assert len(chunks) >= 1
    combined = " ".join(c.text for c in chunks)
    assert "First paragraph" in combined


def test_xlsx_parser(tmp_path):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Age", "City"])
    ws.append(["Alice", 30, "New York"])
    ws.append(["Bob", 25, "London"])
    path = tmp_path / "sample.xlsx"
    wb.save(str(path))
    chunks = parse_document(path)
    assert len(chunks) == 1
    assert "Sheet1" == chunks[0].metadata["sheet"]
    assert "Alice" in chunks[0].text


def test_pptx_parser(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide One"
    slide.placeholders[1].text = "Content of slide one"
    path = tmp_path / "sample.pptx"
    prs.save(str(path))
    chunks = parse_document(path)
    assert len(chunks) >= 1
    assert "Slide One" in chunks[0].text


def test_unsupported_extension(tmp_path):
    f = tmp_path / "binary.exe"
    f.write_bytes(b"\x00\x01\x02")
    with pytest.raises(ValueError, match="Unsupported"):
        parse_document(f)


def test_supported_extensions_set():
    assert ".pdf" in SUPPORTED_EXTENSIONS
    assert ".docx" in SUPPORTED_EXTENSIONS
    assert ".xlsx" in SUPPORTED_EXTENSIONS
    assert ".pptx" in SUPPORTED_EXTENSIONS
    assert ".txt" in SUPPORTED_EXTENSIONS
